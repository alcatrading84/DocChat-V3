"""DocChat Formats v3 — Soporte multi-formato.

Además de PDF, DOCX y TXT, ahora soporta:
- 📊 CSV, Excel (.xlsx)
- 🌐 HTML, Markdown (.md)
- 🖼️ PowerPoint (.pptx)
- 📝 Código fuente (.py, .js, .ts, .java, etc.)
- 📋 JSON, XML, YAML
"""

import os
import csv
import json
import logging
import io
from pathlib import Path
from typing import List, Optional

logger = logging.getLogger(__name__)

# Formatos soportados
SUPPORTED_EXTENSIONS = {
    # Originales
    ".pdf": "PDF Document",
    ".docx": "Word Document",
    ".txt": "Plain Text",
    # Nuevos
    ".md": "Markdown",
    ".html": "HTML",
    ".htm": "HTML",
    ".json": "JSON",
    ".xml": "XML",
    ".yaml": "YAML",
    ".yml": "YAML",
    ".csv": "CSV (Excel/Sheet)",
    ".xlsx": "Excel Workbook",
    ".xls": "Excel Workbook (old)",
    ".pptx": "PowerPoint",
    # Código fuente
    ".py": "Python Code",
    ".js": "JavaScript",
    ".ts": "TypeScript",
    ".java": "Java",
    ".cpp": "C++",
    ".c": "C",
    ".cs": "C#",
    ".go": "Go",
    ".rs": "Rust",
    ".rb": "Ruby",
    ".php": "PHP",
    ".swift": "Swift",
    ".kt": "Kotlin",
    ".sql": "SQL",
    ".sh": "Shell Script",
    ".bat": "Batch Script",
    ".ps1": "PowerShell",
    ".r": "R Script",
}


def get_format_description(ext: str) -> str:
    """Obtener descripción de un formato."""
    return SUPPORTED_EXTENSIONS.get(ext.lower(), "Unknown")


def is_supported(ext: str) -> bool:
    """Verificar si una extensión está soportada."""
    return ext.lower() in SUPPORTED_EXTENSIONS


def load_any_document(filepath: str) -> str:
    """
    Cargar cualquier documento soportado y devolver su texto.

    Args:
        filepath: Ruta al archivo

    Returns:
        Texto extraído del documento

    Raises:
        ValueError: Si el formato no está soportado
    """
    ext = Path(filepath).suffix.lower()
    filename = os.path.basename(filepath)

    # Mapa de cargadores
    loaders = {
        ".txt": _load_txt,
        ".md": _load_txt,  # Markdown es texto plano
        ".py": _load_txt,
        ".js": _load_txt,
        ".ts": _load_txt,
        ".java": _load_txt,
        ".cpp": _load_txt,
        ".c": _load_txt,
        ".cs": _load_txt,
        ".go": _load_txt,
        ".rs": _load_txt,
        ".rb": _load_txt,
        ".php": _load_txt,
        ".swift": _load_txt,
        ".kt": _load_txt,
        ".sql": _load_txt,
        ".sh": _load_txt,
        ".bat": _load_txt,
        ".ps1": _load_txt,
        ".r": _load_txt,
        ".json": _load_json,
        ".xml": _load_xml,
        ".yaml": _load_yaml,
        ".yml": _load_yaml,
        ".html": _load_html,
        ".htm": _load_html,
        ".csv": _load_csv,
        ".xlsx": _load_excel,
        ".xls": _load_excel,
        ".pptx": _load_pptx,
        ".pdf": _load_pdf,
        ".docx": _load_docx,
    }

    loader = loaders.get(ext)
    if not loader:
        raise ValueError(
            f"Formato no soportado: {ext}\n"
            f"Soportados: {', '.join(sorted(SUPPORTED_EXTENSIONS.keys()))}"
        )

    try:
        text = loader(filepath)
        logger.info(f"✅ Cargado: {filename} ({len(text)} caracteres)")
        return text
    except Exception as e:
        raise ValueError(f"Error cargando {filename}: {e}")


# =============================================================================
# CARGADORES INDIVIDUALES
# =============================================================================

def _load_txt(filepath: str) -> str:
    """Cargar archivo de texto plano."""
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _load_json(filepath: str) -> str:
    """Cargar JSON como texto formateado."""
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        data = json.load(f)
    return json.dumps(data, indent=2, ensure_ascii=False)


def _load_xml(filepath: str) -> str:
    """Cargar XML como texto."""
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _load_yaml(filepath: str) -> str:
    """Cargar YAML como texto."""
    try:
        import yaml as yaml_lib
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            data = yaml_lib.safe_load(f)
        return yaml_lib.dump(data, default_flow_style=False, allow_unicode=True)
    except ImportError:
        # Fallback: leer como texto plano
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            return f.read()


def _load_html(filepath: str) -> str:
    """Cargar HTML extrayendo solo el texto."""
    try:
        from html.parser import HTMLParser

        class TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.text = []
                self.skip = False

            def handle_starttag(self, tag, attrs):
                if tag in ("script", "style"):
                    self.skip = True

            def handle_endtag(self, tag):
                if tag in ("script", "style"):
                    self.skip = False
                if tag in ("p", "br", "h1", "h2", "h3", "h4", "h5", "h6", "li"):
                    self.text.append("\n")

            def handle_data(self, data):
                if not self.skip and data.strip():
                    self.text.append(data)

        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            html = f.read()

        parser = TextExtractor()
        parser.feed(html)
        return "\n".join(t.strip() for t in parser.text if t.strip())

    except Exception:
        # Fallback súper simple
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        import re
        # Eliminar etiquetas HTML
        text = re.sub(r"<[^>]+>", " ", text)
        text = re.sub(r"\s+", " ", text)
        return text.strip()


def _load_csv(filepath: str) -> str:
    """Cargar CSV como texto tabular."""
    rows = []
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        for i, row in enumerate(reader):
            if i == 0:
                rows.append(" | ".join(row))
                rows.append("-" * len(rows[0]))
            else:
                rows.append(" | ".join(row))

    return "\n".join(rows)


def _load_excel(filepath: str) -> str:
    """Cargar Excel (.xlsx/.xls) como texto."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ImportError(
            "Para leer Excel necesitas: pip install openpyxl"
        )

    wb = load_workbook(filepath, read_only=True, data_only=True)
    parts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"\n=== Hoja: {sheet_name} ===\n")

        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = " | ".join(cells)
            if line.strip():
                parts.append(line)

    wb.close()
    return "\n".join(parts)


def _load_pptx(filepath: str) -> str:
    """Cargar PowerPoint (.pptx) como texto."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "Para leer PowerPoint necesitas: pip install python-pptx"
        )

    prs = Presentation(filepath)
    parts = []

    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"\n--- Diapositiva {i} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    cells = [cell.text for cell in row.cells]
                    parts.append(" | ".join(cells))

    return "\n".join(parts)


def _load_pdf(filepath: str) -> str:
    """Cargar PDF (igual que en engine.py original)."""
    from pypdf import PdfReader
    try:
        reader = PdfReader(filepath)
    except Exception as e:
        raise ValueError(f"No se pudo abrir el PDF: {e}")

    if reader.is_encrypted:
        try:
            reader.decrypt("")
        except Exception:
            raise ValueError(
                "PDF protegido con contraseña. "
                "Guarda una copia sin protección e intenta de nuevo."
            )

    text = []
    for page in reader.pages:
        try:
            t = page.extract_text()
            if t and t.strip():
                text.append(t)
        except Exception:
            pass

    return "\n".join(text)


def _load_docx(filepath: str) -> str:
    """Cargar Word (.docx)."""
    from docx import Document as DocxDoc
    doc = DocxDoc(filepath)
    return "\n".join(p.text for p in doc.paragraphs)
